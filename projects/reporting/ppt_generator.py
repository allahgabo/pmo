"""
PowerPoint Report Generator
Generates executive PMO reports in PowerPoint format
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime
from django.db.models import Count, Avg, Sum, Q
from projects.models import Project, Risk, Task


class PowerPointGenerator:
    """Generate PMO reports in PowerPoint format"""
    
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        
    def create_portfolio_report(self, output_path=None):
        """Create complete portfolio status report"""
        # Title slide
        self._add_title_slide()
        
        # Executive summary
        self._add_executive_summary()
        
        # Portfolio health
        self._add_portfolio_health()
        
        # Project status
        self._add_project_status()
        
        # Risk summary
        self._add_risk_summary()
        
        # Budget status
        self._add_budget_status()
        
        # Recommendations
        self._add_recommendations()
        
        # Save
        if not output_path:
            output_path = f'PMO_Report_{datetime.now().strftime("%Y%m%d")}.pptx'
        
        self.prs.save(output_path)
        return output_path
    
    def _add_title_slide(self):
        """Add title slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "PMO Portfolio Status Report"
        subtitle.text = f"Generated: {datetime.now().strftime('%B %d, %Y')}"
        
        # Style title
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)
    
    def _add_executive_summary(self):
        """Add executive summary slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Executive Summary"
        
        # Get portfolio stats
        projects = Project.objects.all()
        total_projects = projects.count()
        on_track = projects.filter(status='on_track').count()
        at_risk = projects.filter(status='at_risk').count()
        delayed = projects.filter(status='delayed').count()
        
        total_budget = projects.aggregate(Sum('budget'))['budget__sum'] or 0
        total_spent = projects.aggregate(Sum('spent'))['spent__sum'] or 0
        
        avg_health = sum([p.health_score for p in projects]) / total_projects if total_projects else 0
        
        # Add content
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        
        tf.text = f"Portfolio Overview ({total_projects} Projects)"
        p = tf.add_paragraph()
        p.text = f"\n✓ On Track: {on_track} projects ({on_track/total_projects*100:.0f}%)"
        p.font.size = Pt(18)
        
        p = tf.add_paragraph()
        p.text = f"⚠ At Risk: {at_risk} projects ({at_risk/total_projects*100:.0f}%)"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(245, 158, 11)
        
        p = tf.add_paragraph()
        p.text = f"✗ Delayed: {delayed} projects ({delayed/total_projects*100:.0f}%)"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(239, 68, 68)
        
        p = tf.add_paragraph()
        p.text = f"\n💰 Total Budget: ${total_budget:,.0f}"
        p.font.size = Pt(18)
        
        p = tf.add_paragraph()
        p.text = f"💸 Total Spent: ${total_spent:,.0f} ({total_spent/total_budget*100:.0f}%)"
        p.font.size = Pt(18)
        
        p = tf.add_paragraph()
        p.text = f"\n📊 Average Health Score: {avg_health:.0f}/100"
        p.font.size = Pt(18)
    
    def _add_portfolio_health(self):
        """Add portfolio health slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Portfolio Health Dashboard"
        
        # Add metrics as text boxes
        projects = Project.objects.all()
        
        metrics = [
            ('Projects', projects.count()),
            ('Avg Completion', f"{sum([p.completion_percentage for p in projects])/projects.count():.1f}%"),
            ('Avg SPI', f"{sum([float(p.spi) for p in projects])/projects.count():.2f}"),
            ('Avg CPI', f"{sum([float(p.cpi) for p in projects])/projects.count():.2f}"),
        ]
        
        # Add metric boxes
        box_width = Inches(2)
        box_height = Inches(1.5)
        start_left = Inches(0.5)
        top = Inches(2)
        spacing = Inches(2.5)
        
        for i, (label, value) in enumerate(metrics):
            left = start_left + (i * spacing)
            
            # Add box
            shape = slide.shapes.add_shape(
                1,  # Rectangle
                left, top, box_width, box_height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(102, 126, 234)
            
            # Add text
            tf = shape.text_frame
            tf.clear()
            
            p = tf.paragraphs[0]
            p.text = str(value)
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
            
            p = tf.add_paragraph()
            p.text = label
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
    
    def _add_project_status(self):
        """Add project status details slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Project Status Detail"
        
        projects = Project.objects.all()[:5]  # Top 5 projects
        
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        
        for project in projects:
            status_icon = {
                'on_track': '✓',
                'at_risk': '⚠',
                'delayed': '✗',
                'completed': '✓'
            }.get(project.status, '•')
            
            status_color = {
                'on_track': RGBColor(34, 197, 94),
                'at_risk': RGBColor(245, 158, 11),
                'delayed': RGBColor(239, 68, 68),
                'completed': RGBColor(34, 197, 94)
            }.get(project.status, RGBColor(0, 0, 0))
            
            p = tf.add_paragraph()
            p.text = f"{status_icon} {project.name}"
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = status_color
            
            p = tf.add_paragraph()
            p.text = f"   Completion: {project.completion_percentage}% | Health: {project.health_score}/100 | SPI: {project.spi}"
            p.font.size = Pt(14)
            p.level = 1
    
    def _add_risk_summary(self):
        """Add risk summary slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Top Risks"
        
        risks = Risk.objects.filter(status__in=['open', 'mitigating']).order_by('-severity', '-probability')[:10]
        
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        
        for risk in risks:
            severity_icon = {
                'critical': '🔴',
                'high': '🟠',
                'medium': '🟡',
                'low': '🟢'
            }.get(risk.severity, '⚪')
            
            p = tf.add_paragraph()
            p.text = f"{severity_icon} [{risk.project.code}] {risk.title}"
            p.font.size = Pt(14)
            
            p = tf.add_paragraph()
            p.text = f"   Probability: {risk.probability}% | Impact: {risk.impact}/10 | Score: {risk.risk_score:.2f}"
            p.font.size = Pt(12)
            p.level = 1
    
    def _add_budget_status(self):
        """Add budget status slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Budget Status"
        
        projects = Project.objects.all()
        
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        
        for project in projects:
            budget_pct = (float(project.spent) / float(project.budget) * 100) if project.budget else 0
            
            status = '✓' if budget_pct < 100 else '✗'
            color = RGBColor(34, 197, 94) if budget_pct < 100 else RGBColor(239, 68, 68)
            
            p = tf.add_paragraph()
            p.text = f"{status} {project.name}"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = color
            
            p = tf.add_paragraph()
            p.text = f"   Budget: ${project.budget:,.0f} | Spent: ${project.spent:,.0f} ({budget_pct:.1f}%)"
            p.font.size = Pt(12)
            p.level = 1
    
    def _add_recommendations(self):
        """Add recommendations slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Recommendations"
        
        # Generate AI-powered recommendations
        projects = Project.objects.all()
        at_risk_projects = projects.filter(Q(status='at_risk') | Q(status='delayed'))
        
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        
        recommendations = []
        
        if at_risk_projects.exists():
            recommendations.append(f"• {at_risk_projects.count()} projects require immediate attention")
        
        low_spi_projects = [p for p in projects if float(p.spi) < 0.9]
        if low_spi_projects:
            recommendations.append(f"• {len(low_spi_projects)} projects are behind schedule - consider resource reallocation")
        
        over_budget = [p for p in projects if float(p.spent) > float(p.budget)]
        if over_budget:
            recommendations.append(f"• {len(over_budget)} projects over budget - conduct cost reviews")
        
        high_risks = Risk.objects.filter(severity__in=['high', 'critical'], status='open').count()
        if high_risks > 0:
            recommendations.append(f"• {high_risks} high/critical risks need mitigation plans")
        
        if not recommendations:
            recommendations.append("• Portfolio is healthy - maintain current monitoring")
        
        for rec in recommendations:
            p = tf.add_paragraph()
            p.text = rec
            p.font.size = Pt(18)
            p.space_after = Pt(12)


def generate_portfolio_report(output_path=None):
    """Convenience function to generate portfolio report"""
    generator = PowerPointGenerator()
    return generator.create_portfolio_report(output_path)